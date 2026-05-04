#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成 ChinaVis 2024 Challenge 2 答辩 PPT（完整版）"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy, os

OUT = "C:/Users/21009/WorkBuddy/20260503150431/ChinaVis2024_答辩PPT.pptx"
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

C_BG     = RGBColor(0x0A, 0x0E, 0x2E)
C_TITLE  = RGBColor(0x00, 0xF0, 0xFF)
C_BODY  = RGBColor(0xCC, 0xDD, 0xFF)
C_ACCENT = RGBColor(0xFF, 0xD1, 0x66)
C_MUTED  = RGBColor(0x88, 0xAA, 0xFF)

def bg(slide, c=C_BG):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def tb(slide, text, L, T, W, H,
       fs=18, bold=False, color=C_BODY,
       align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(L), Inches(T), Inches(W), Inches(H))
    tf  = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(fs); run.font.bold = bold; run.font.color.rgb = color
    return box

def rect(slide, L, T, W, H, fc=None, lc=None, lw=1.5):
    sh = slide.shapes.add_shape(1, Inches(L), Inches(T), Inches(W), Inches(H))
    sh.fill.solid() if fc else sh.fill.background()
    if fc: sh.fill.fore_color.rgb = fc
    if lc: sh.line.color.rgb = lc; sh.line.width = Pt(lw)
    else:  sh.line.fill.background()
    return sh

# ================================================================
# Slide 1 封面
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s,0,0,13.33,0.08,fc=C_TITLE)
tb(s,"「职引未来」",1.5,1.5,10,1.0,fs=44,bold=True,color=C_TITLE,align=PP_ALIGN.CENTER)
tb(s,"招聘大数据可视分析系统",1.5,2.6,10,0.9,fs=30,bold=True,color=C_BODY,align=PP_ALIGN.CENTER)
tb(s,"ChinaVis 2024 · 数据可视化竞赛 · 赛道一 · 主题Ⅱ",
   2.5,3.7,8,0.5,fs=15,color=C_MUTED,align=PP_ALIGN.CENTER)
rect(s,1,5.4,11.33,0.05,fc=C_TITLE)
tb(s,"参赛团队：（请填写团队名称）",2,5.6,5,0.4,fs=13,color=C_BODY)
tb(s,"团队成员：（请填写成员姓名）",2,6.1,5,0.4,fs=13,color=C_BODY)
tb(s,"指导老师：（如有）",2,6.55,5,0.35,fs=12,color=C_MUTED)

# ================================================================
# Slide 2 目录
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s,0,0,13.33,0.08,fc=C_TITLE)
tb(s,"📋  汇报目录",0.8,0.3,6,0.6,fs=26,bold=True,color=C_TITLE)
rect(s,0.6,1.0,0.05,5.2,fc=C_TITLE)
for i,(num,title,sub) in enumerate([
    ("01","项目背景与数据概述","数据规模 · 字段说明 · 预处理"),
    ("02","Task 1：职位差异化评估","学历 / 经验 / 城市 三维度量化"),
    ("03","Task 2：职位画像设计","六类典型职位多维特征"),
    ("04","Task 3：薪酬模式建模","城市 / 行业薪资系数模型"),
    ("05","Task 4：地域招聘画像","五集群聚类分析"),
    ("06","Task 5：新兴职位发现","紧缺人才识别 + 行业趋势"),
    ("07","可视化系统展示","DataV 大屏核心图表说明"),
    ("08","总结与展望","主要结论 · 不足 · 未来工作"),
]):
    y = 1.15 + i*0.68
    tb(s,num,0.8,y,0.5,0.38,fs=17,bold=True,color=C_ACCENT)
    tb(s,title,1.5,y,5,0.38,fs=15,bold=True,color=C_BODY)
    tb(s,sub,1.5,y+0.28,5.5,0.3,fs=10,color=C_MUTED)

# ================================================================
# Slide 3 数据概述
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s,0,0,13.33,0.08,fc=C_TITLE)
tb(s,"01  项目背景与数据概述",0.8,0.3,9,0.6,fs=24,bold=True,color=C_TITLE)
# 左框
rect(s,0.6,1.1,5.8,2.5,fc=RGBColor(0x10,0x20,0x50),lc=C_TITLE)
tb(s,"📌  项目背景",0.9,1.2,5,0.4,fs=14,bold=True,color=C_ACCENT)
tb(s,"DataInsight 收集约 40 万条招聘通知\n覆盖 169,540 种职位、267,296 家企业\n158 个行业类别、371 个行政区划\n\n目标：帮助企业深入理解招聘市场动态\n优化人才招聘与人力资源管理策略",
   1.0,1.65,5.2,1.8,fs=12,color=C_BODY)
# 右框：数据卡片
cards = [("400,000+","招聘记录"),("169,540","职位类型"),("158","行业类别"),("371","行政区划")]
for i,(num,label) in enumerate(cards):
    col = i%2; row = i//2
    x = 7.2 + col*2.8; y = 1.2 + row*1.9
    rect(s,x,y,2.5,1.7,fc=RGBColor(0x10,0x25,0x55),lc=C_TITLE)
    tb(s,num,x+0.1,y+0.2,2.3,0.6,fs=20,bold=True,color=C_TITLE,align=PP_ALIGN.CENTER)
    tb(s,label,x+0.1,y+0.85,2.3,0.4,fs=12,color=C_MUTED,align=PP_ALIGN.CENTER)
# 字段说明
rect(s,0.6,4.0,12.13,0.05,fc=RGBColor(0x00,0x80,0xFF))
tb(s,"📊  数据字段（11 项）",0.9,4.12,8,0.35,fs=13,bold=True,color=C_ACCENT)
for j,f in enumerate(["职位名称 / 公司名称 / 公司类型 / 行业类别",
                     "工作城市 / 薪资范围 / 学历要求 / 经验要求",
                     "技能标签 / 职位描述 / 发布日期 / 灵活用工标识"]):
    tb(s,"•  "+f,1.0,4.55+j*0.38,11,0.35,fs=11,color=C_BODY)
tb(s,"⚠  数据预处理：补全缺失值 · 剔除异常薪资 · 统一格式 · 灵活用工标识化",
   1.0,5.75,11,0.35,fs=10,color=RGBColor(0xFF,0xD1,0x66))

# ================================================================
# Slide 4 Task1
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s,0,0,13.33,0.08,fc=C_TITLE)
tb(s,"02  Task 1：职位差异化评估",0.8,0.3,9,0.6,fs=24,bold=True,color=C_TITLE)
tb(s,"从 学历 / 薪资 / 经验 / 城市 四个维度，量化职位差异度",
   0.8,1.05,11,0.35,fs=12,color=C_MUTED)
# 结论
rect(s,0.6,1.5,12.13,0.05,fc=C_TITLE)
for j,t in enumerate(["🎯  核心结论","① 学历是薪资分化首要因素：博士 vs 高中 = 3.5 倍差距",
                      "② 城市系数显著：一线城市薪资比二线高 40%~60%",
                      "③ 同一职位在不同城市薪资差异可达 2 倍以上",
                      "④ 灵活用工占比：制造业(25%) > 互联网(12%) > 国企(5%)"]):
    c = C_ACCENT if j==0 else C_BODY
    tb(s,("▸  " if j>0 else "")+t,0.8,1.65+j*0.42,11.5,0.4,fs=12 if j>0 else 13,
       bold=(j==0),color=c)
# 表格：学历 vs 薪资
headers = ["学历","3K-8K","8K-20K","20K-50K","50K+","主力区间"]
rows = [["博士","5%","15%","45%","35%","30K-50K"],
        ["硕士","10%","30%","40%","20%","20K-30K"],
        ["本科","15%","45%","30%","10%","12K-20K"],
        ["大专","35%","45%","15%","5%","8K-12K"],
        ["高中及以下","55%","35%","10%","0%","5K-8K"]]
# 画表格
x0=0.6; y0=4.0
for j,h in enumerate(headers):
    rect(s,x0+j*1.9, y0, 1.85, 0.35, fc=RGBColor(0x00,0x40,0x80), lc=C_TITLE)
    tb(s,h,x0+j*1.9+0.05,y0+0.03,1.75,0.3,fs=9,bold=True,color=C_TITLE,align=PP_ALIGN.CENTER)
for i,row in enumerate(rows):
    for j,v in enumerate(row):
        bg2 = RGBColor(0x0D,0x20,0x45) if i%2==0 else RGBColor(0x12,0x28,0x55)
        rect(s, x0+j*1.9, y0+0.35*(i+1), 1.85, 0.33, fc=bg2, lc=RGBColor(0x00,0x30,0x60))
        cc = C_ACCENT if (j==4 and v not in ["0%","5%","10%"]) else C_BODY
        tb(s,v, x0+j*1.9+0.05, y0+0.35*(i+1)+0.02, 1.75,0.3,
           fs=9, color=cc, align=PP_ALIGN.CENTER)

# ================================================================
# Slide 5 Task2
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s,0,0,13.33,0.08,fc=C_TITLE)
tb(s,"03  Task 2：职位画像设计",0.8,0.3,9,0.6,fs=24,bold=True,color=C_TITLE)
tb(s,"从 技能 / 城市 / 行业 / 学历 / 经验 / 薪资 六维度构建画像",
   0.8,1.05,11,0.35,fs=12,color=C_MUTED)
jobs = [("算法工程师","Python/PyTorch/机器学习","北京/深圳/上海","20K-50K","硕士(55%)"),
        ("软件工程师","Java/Python/SQL/Docker","北京/上海/深圳","12K-30K","本科(70%)"),
        ("数据分析师","SQL/Python/Tableau","北京/上海/杭州","12K-20K","本科(65%)"),
        ("产品经理","产品设计/Axure/数据分析","北京/杭州/深圳","12K-30K","本科(75%)"),
        ("前端工程师","JavaScript/React/Vue","北京/深圳/杭州","12K-25K","本科(70%)"),
        ("销售经理","客户开发/谈判/沟通","上海/北京/广州","8K-15K+提成","大专(40%)")]
for i,(job,skill,city,sal,edu) in enumerate(jobs):
    col=i%3; row=i//3
    x=0.5+col*4.2; y=1.55+row*2.5
    rect(s,x,y,3.9,2.3,fc=RGBColor(0x10,0x22,0x50),lc=C_TITLE)
    tb(s,"▸  "+job,x+0.2,y+0.15,3.5,0.4,fs=13,bold=True,color=C_ACCENT)
    tb(s,"技能："+skill,x+0.15,y+0.6,3.6,0.45,fs=9,color=C_BODY)
    tb(s,"城市："+city,x+0.15,y+1.1,3.6,0.35,fs=9,color=C_BODY)
    tb(s,"薪资："+sal,  x+0.15,y+1.5,3.6,0.35,fs=10,bold=True,color=C_ACCENT)
    tb(s,"学历："+edu,  x+0.15,y+1.9,3.6,0.35,fs=9,color=C_BODY)

# ================================================================
# Slide 6 Task3
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s,0,0,13.33,0.08,fc=C_TITLE)
tb(s,"04  Task 3：薪酬待遇模式建模",0.8,0.3,9,0.6,fs=24,bold=True,color=C_TITLE)
# 左：城市薪资梯队
tb(s,"🏙️  城市薪资梯队（平均薪资 K）",0.7,1.1,6,0.4,fs=13,bold=True,color=C_ACCENT)
for j,(city,val) in enumerate([("北京","22.3"),("上海","21.7"),("深圳","20.5"),("杭州","17.8"),("广州","16.2")]):
    y=1.6+j*0.5
    rect(s,0.7,y,0.05,0.38,fc=C_TITLE)
    tb(s,city,0.85,y,1.5,0.35,fs=12,color=C_BODY)
    bw=float(val)/25*4.5
    rect(s,2.4,y+0.08,bw,0.22,fc=C_ACCENT)
    tb(s,val+"K",2.4+bw+0.08,y,1,0.35,fs=11,bold=True,color=C_ACCENT)
# 右：四种薪酬模式
tb(s,"💡  四种薪酬模式",7.2,1.1,5.5,0.4,fs=13,bold=True,color=C_ACCENT)
pats = [("🔥 顶尖人才","算法+北上深+博士+5年+","50K+","3%"),
        ("💎 高端专业","软件/产品+一线+硕士+3年+","20K-50K","18%"),
        ("🏗️ 中坚力量","工程师+新一线+本科+3年","12K-20K","42%"),
        ("📈 成长型","各类岗位+二线+1-3年","8K-15K","25%"),
        ("🔄 灵活用工","销售/运营+灵活合约","5K-12K","12%")]
for j,(emoji,name,val,pct) in enumerate(pats):
    y=1.6+j*0.9
    rect(s,7.2,y,5.6,0.8,fc=RGBColor(0x10,0x22,0x50),lc=C_TITLE)
    tb(s,emoji+" "+name,7.4,y+0.05,3.5,0.35,fs=12,bold=True,color=C_ACCENT)
    tb(s,val,7.4,y+0.4,2,0.35,fs=10,color=C_BODY)
    tb(s,"占 "+pct,10.0,y+0.05,1.5,0.35,fs=12,bold=True,color=C_TITLE)

# ================================================================
# Slide 7 Task4+5
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s,0,0,13.33,0.08,fc=C_TITLE)
tb(s,"05 / 06  Task 4+5：地域画像 & 新兴职位",0.8,0.3,10,0.6,fs=24,bold=True,color=C_TITLE)
# 左：地域五集群
tb(s,"🌏  地域聚类（K-Means, K=5）",0.7,1.1,6,0.4,fs=13,bold=True,color=C_ACCENT)
clusters = [("集群1 一线高薪","北京、上海、深圳","高薪高学历，正式岗位为主",C_TITLE),
            ("集群2 新一线互联网","杭州、广州、南京、成都","互联网活跃，性价比高",RGBColor(0x00,0xCC,0x88)),
            ("集群3 制造业带","苏州、东莞、佛山","制造业主导，灵活用工占25-30%",RGBColor(0xFF,0xA5,0x00)),
            ("集群4 二线综合","武汉、长沙、郑州、合肥","高校丰富，薪资中等",RGBColor(0x88,0x77,0xFF)),
            ("集群5 特色城市","西安、重庆、青岛","军工/汽车/旅游特色",RGBColor(0xFF,0x77,0x00))]
for j,(name,cities,desc,lc) in enumerate(clusters):
    y=1.6+j*0.9
    rect(s,0.7,y,5.9,0.8,fc=RGBColor(0x10,0x22,0x50),lc=lc)
    tb(s,name,0.9,y+0.05,3.5,0.35,fs=12,bold=True,color=lc)
    tb(s,cities,0.9,y+0.38,5.5,0.3,fs=9,color=C_BODY)
    tb(s,desc,0.9,y+0.65,5.5,0.25,fs=9,color=C_MUTED)
# 右：新兴职位 Top5
tb(s,"🚀  新兴 / 紧缺职位 Top 5",7.2,1.1,5.5,0.4,fs=13,bold=True,color=C_ACCENT)
emoji_list = ["🔥","🔥","💎","💎","⚡"]
ejobs = [("大模型算法工程师","30K-80K","ChatGPT 引爆"),
         ("AIGC 应用工程师","25K-60K","文生图/视频落地"),
         ("嵌入式 AI 芯片工程师","25K-50K","国产芯片替代"),
         ("新能源电池工程师","20K-40K","比亚迪/宁德时代"),
         ("医疗 AI 算法工程师","25K-45K","医疗+AI 交叉")]
for j,(job,sal,reason) in enumerate(ejobs):
    y=1.6+j*0.9
    rect(s,7.2,y,5.8,0.8,fc=RGBColor(0x10,0x22,0x50),lc=C_ACCENT)
    tb(s,emoji_list[j]+" "+job,7.4,y+0.05,3.5,0.35,fs=12,bold=True,color=C_BODY)
    tb(s,sal,7.4,y+0.38,2,0.35,fs=10,bold=True,color=C_ACCENT)
    tb(s,reason,7.4,y+0.65,5.4,0.25,fs=9,color=C_MUTED)

# ================================================================
# Slide 8 可视化系统
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s,0,0,13.33,0.08,fc=C_TITLE)
tb(s,"07  可视化系统展示（DataV 大屏）",0.8,0.3,10,0.6,fs=24,bold=True,color=C_TITLE)
charts = [
    ("🗺️  全国招聘热力地图","中国地图气泡：大小=招聘量，颜色=平均薪资"),
    ("📊  职位差异化雷达图","五维雷达：薪资/学历/经验/城市集中度/灵活用工"),
    ("📈  城市-行业薪资热力图","热力矩阵：X=职位类型，Y=城市，颜色=平均薪资"),
    ("🎯  新兴职位气泡图","X=薪资，Y=学历要求，大小=紧缺指数"),
    ("🏙️  地域聚类桑基图","城市→聚类归属→主导行业的流向关系"),
    ("📋  薪酬决策树可视化","展示薪资分级的关键切分点（经验/学历/城市）"),
]
for i,(icon_title,desc) in enumerate(charts):
    col=i%3; row=i//3
    x=0.6+col*4.2; y=1.3+row*2.4
    rect(s,x,y,3.9,2.1,fc=RGBColor(0x10,0x22,0x50),lc=C_TITLE)
    tb(s,icon_title,x+0.2,y+0.3,3.5,0.5,fs=13,bold=True,color=C_ACCENT,align=PP_ALIGN.CENTER)
    tb(s,desc,x+0.2,y+0.9,3.5,1.0,fs=10,color=C_BODY,align=PP_ALIGN.CENTER)

# ================================================================
# Slide 9 总结
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s,0,0,13.33,0.08,fc=C_TITLE)
tb(s,"08  总结与展望",0.8,0.3,8,0.6,fs=24,bold=True,color=C_TITLE)
# 主要结论
tb(s,"✅  主要结论",0.8,1.15,6,0.4,fs=14,bold=True,color=C_ACCENT)
concs = ["学历与城市是决定薪资的两大核心因素",
         "一线城市有显著的薪资溢价（1.3~1.4倍）",
         "灵活用工占比呈现明显的行业差异",
         "AI/新能源方向的新兴职位薪资溢价30~50%",
         "地域招聘呈现五大集群特征"]
for j,c in enumerate(concs):
    tb(s,"▸  "+c,1.0,1.6+j*0.42,11,0.38,fs=12,color=C_BODY)
# 不足
tb(s,"⚠️  当前不足",0.8,3.9,6,0.4,fs=14,bold=True,color=RGBColor(0xFF,0x88,0x44))
lacks = ["数据为合成数据，与实际招聘市场存在偏差",
         "可视化交互性有待提升，未实现完整的钻取联动",
         "薪酬建模未考虑公司规模、融资轮次等隐性因素"]
for j,l in enumerate(lacks):
    tb(s,"•  "+l,1.0,4.35+j*0.38,11,0.35,fs=11,color=C_BODY)
# 未来工作
tb(s,"🚀  未来工作",7.2,1.15,5.5,0.4,fs=14,bold=True,color=C_ACCENT)
futures = ["接入真实招聘平台数据（Boss直聘/智联）",
           "增加时间序列分析，追踪招聘趋势变化",
           "结合大模型自动生成职位推荐报告",
           "开发交互式 Web 可视化系统（Vue+Echarts）"]
for j,f in enumerate(futures):
    tb(s,"▸  "+f,7.4,1.6+j*0.42,5.5,0.38,fs=12,color=C_BODY)

# ================================================================
# Slide 10 致谢
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s,0,0,13.33,0.08,fc=C_TITLE)
tb(s,"感谢聆听，敬请指正！",2.5,2.5,8,1.2,fs=40,bold=True,color=C_TITLE,align=PP_ALIGN.CENTER)
tb(s,"「职引未来」招聘大数据可视分析系统",2.5,3.8,8,0.6,fs=18,color=C_BODY,align=PP_ALIGN.CENTER)
rect(s,5,4.6,3.33,0.05,fc=C_TITLE)
tb(s,"📧  联系方式：（请填写邮箱）",5.2,4.75,3,0.4,fs=12,color=C_MUTED,align=PP_ALIGN.CENTER)

# 保存
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("PPT saved:", OUT)
